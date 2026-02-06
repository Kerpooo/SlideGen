from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import HTMLResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from fastapi import Request
from pathlib import Path
import tempfile
import shutil
from typing import List
from pptx import Presentation
from copy import deepcopy
import os
import io
import hashlib

app = FastAPI(title="SlideGen - PPTX Batch Processor")

# Configurar templates
templates = Jinja2Templates(directory="templates")

# Montar archivos estáticos
app.mount("/static", StaticFiles(directory="static"), name="static")


def static_url(filename: str) -> str:
    """Genera URL estática con hash para cache busting."""
    file_path = Path("static") / filename
    if file_path.exists():
        content_hash = hashlib.md5(file_path.read_bytes()).hexdigest()[:8]
        return f"/static/{filename}?v={content_hash}"
    return f"/static/{filename}"


templates.env.globals["static_url"] = static_url


@app.get("/", response_class=HTMLResponse)
async def home(request: Request):
    """Renderiza la página principal."""
    return templates.TemplateResponse("index.html", {"request": request})


@app.post("/api/process")
async def process_slides(
    template: UploadFile = File(...),
    names: str = Form(...),
    export_pdf: bool = Form(False),
    email_results: bool = Form(False)
):
    """
    Procesa el archivo PPTX con la lista de nombres y devuelve el archivo directamente.
    """
    # Validar archivo
    if not template.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(status_code=400, detail="Solo se aceptan archivos .pptx o .ppt")

    # Procesar lista de nombres
    names_list = [name.strip() for name in names.split('\n') if name.strip()]

    if not names_list:
        raise HTTPException(status_code=400, detail="Debe proporcionar al menos un nombre")

    # Guardar archivo temporal de entrada
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_input:
        shutil.copyfileobj(template.file, tmp_input)
        tmp_input_path = tmp_input.name

    # Archivo temporal de salida
    tmp_output_path = tempfile.mktemp(suffix=".pptx")

    try:
        # Procesar la presentación
        process_presentation(tmp_input_path, names_list, tmp_output_path)

        # Leer el archivo procesado en memoria
        with open(tmp_output_path, "rb") as f:
            file_content = io.BytesIO(f.read())

        # Generar nombre del archivo de salida
        original_name = template.filename.rsplit('.', 1)[0]
        output_filename = f"{original_name}_procesado.pptx"

        return StreamingResponse(
            file_content,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{output_filename}"'
            }
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error procesando archivo: {str(e)}")

    finally:
        # Limpiar archivos temporales
        if os.path.exists(tmp_input_path):
            os.unlink(tmp_input_path)
        if os.path.exists(tmp_output_path):
            os.unlink(tmp_output_path)


def process_presentation(template_path: str, names: List[str], output_path: str):
    """
    Procesa la presentación reemplazando el marcador {{NOMBRE}} con cada nombre.
    Las slides generadas se insertan justo después de la plantilla original.
    """
    from generar_slides import (
        get_template_slide_indices,
        duplicate_slide_with_images,
        replace_marker_in_slide,
        delete_slide,
        move_slide
    )

    MARKER = "{{NOMBRE}}"

    prs = Presentation(template_path)

    # Encontrar todas las diapositivas plantilla
    template_indices = get_template_slide_indices(prs, MARKER)

    if not template_indices:
        raise ValueError(f"No se encontró el marcador {MARKER} en ninguna diapositiva")

    # Procesar cada plantilla y sus nombres
    # Llevamos cuenta del offset porque al insertar slides los índices cambian
    offset = 0

    for template_idx in template_indices:
        # Posición actual de la plantilla (considerando inserciones previas)
        current_template_pos = template_idx + offset
        template_slide = prs.slides[current_template_pos]

        # Generar copias para cada nombre
        for i, name in enumerate(names):
            # Duplicar la slide (se agrega al final)
            new_slide = duplicate_slide_with_images(prs, template_slide)
            replace_marker_in_slide(new_slide, name, MARKER)

            # Mover la slide desde el final a justo después de la plantilla
            last_index = len(prs.slides) - 1
            insert_position = current_template_pos + 1 + i
            move_slide(prs, last_index, insert_position)

        # Actualizar offset para la siguiente plantilla
        offset += len(names)

    # Eliminar las diapositivas plantilla originales (ahora están en posiciones diferentes)
    # Calcular las nuevas posiciones de las plantillas
    templates_to_delete = []
    current_offset = 0
    for original_idx in template_indices:
        new_pos = original_idx + current_offset
        templates_to_delete.append(new_pos)
        current_offset += len(names)

    # Eliminar de mayor a menor para no afectar índices
    for i in sorted(templates_to_delete, reverse=True):
        delete_slide(prs, i)

    prs.save(output_path)