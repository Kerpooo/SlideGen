from pptx import Presentation
from copy import deepcopy

# ---------------- CONFIG ----------------
TEMPLATE_PPTX = "plantilla.pptx"
OUTPUT_PPTX = "resultado_con_nombres.pptx"
NAMES_FILE = "nombres.txt"

MARKER = "{{NOMBRE}}"
# ---------------------------------------


def slide_has_marker(slide, marker):
    """Verifica si una diapositiva contiene el marcador."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in p.runs)
                if marker in full_text:
                    return True
    return False


def get_template_slide_indices(prs, marker):
    """Retorna los índices de todas las diapositivas que contienen el marcador."""
    indices = []
    for i, slide in enumerate(prs.slides):
        if slide_has_marker(slide, marker):
            indices.append(i)
    return indices


def duplicate_slide_with_images(prs, source_slide):
    """
    Duplica una diapositiva copiando shapes Y relaciones (imágenes, etc).
    Mapea los rId originales a los nuevos para que las imágenes funcionen.
    """
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Limpiar shapes que vienen del layout
    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Mapeo de rId antiguo → rId nuevo
    rid_map = {}

    # Copiar relaciones (imágenes, etc) de la slide original
    for rel in source_slide.part.rels.values():
        # Evitar duplicar relaciones al slide layout/notes
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue

        # Si es una imagen u otro recurso embebido, copiar la relación
        if rel.is_external:
            new_rid = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            try:
                new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
            except Exception:
                continue

        # Guardar el mapeo de rId antiguo a nuevo
        rid_map[rel.rId] = new_rid

    # Copiar los elementos XML de los shapes y actualizar rIds
    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)

        # Actualizar todos los rId en el elemento copiado
        for elem in new_el.iter():
            for attr_name in list(elem.attrib.keys()):
                # Buscar atributos que terminan en 'embed', 'link', 'id' (referencias a relaciones)
                if attr_name.endswith('}embed') or attr_name.endswith('}link') or attr_name.endswith('}id'):
                    old_rid = elem.attrib[attr_name]
                    if old_rid in rid_map:
                        elem.attrib[attr_name] = rid_map[old_rid]

        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Copiar background si existe
    try:
        if source_slide.background.fill.type is not None:
            bg_el = deepcopy(source_slide.background._element)
            # Actualizar rIds en el background también
            for elem in bg_el.iter():
                for attr_name in list(elem.attrib.keys()):
                    if attr_name.endswith('}embed') or attr_name.endswith('}link'):
                        old_rid = elem.attrib[attr_name]
                        if old_rid in rid_map:
                            elem.attrib[attr_name] = rid_map[old_rid]
            new_slide.background._element = bg_el
    except Exception:
        pass

    return new_slide


def replace_marker_runs(shape, name, marker=None):
    """
    Reemplaza MARKER dentro de los runs existentes.
    Preserva estilo y posición.
    """
    if marker is None:
        marker = MARKER
    
    tf = shape.text_frame
    changed = False

    for p in tf.paragraphs:
        for run in p.runs:
            if marker in run.text:
                run.text = run.text.replace(marker, name)
                changed = True

    return changed


def replace_marker_robust(shape, name, marker=None):
    """
    Fallback: si el marcador está partido en varios runs,
    reconstruye el texto manteniendo el estilo del primer run.
    """
    if marker is None:
        marker = MARKER
    
    tf = shape.text_frame
    changed = False

    for p in tf.paragraphs:
        full = "".join(run.text for run in p.runs)
        if marker not in full:
            continue

        new_text = full.replace(marker, name)

        if p.runs:
            r0 = p.runs[0]
            for r in list(p.runs)[1:]:
                r._r.getparent().remove(r._r)
            r0.text = new_text
        else:
            p.text = new_text

        changed = True

    return changed


def replace_marker_in_slide(slide, name, marker=None):
    """Reemplaza el marcador en todos los shapes de la slide."""
    if marker is None:
        marker = MARKER
    
    found_any = False

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if replace_marker_runs(shape, name, marker):
            found_any = True

    if not found_any:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if replace_marker_robust(shape, name, marker):
                found_any = True

    return found_any


def delete_slide(prs, index):
    """Elimina una diapositiva por índice."""
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[index]


def move_slide(prs, old_index, new_index):
    """Mueve una diapositiva de una posición a otra."""
    slides = prs.slides._sldIdLst
    slide = slides[old_index]
    del slides[old_index]
    slides.insert(new_index, slide)


def main():
    prs = Presentation(TEMPLATE_PPTX)

    # Leer nombres
    with open(NAMES_FILE, "r", encoding="utf-8") as f:
        names = [line.strip() for line in f if line.strip()]

    # Encontrar todas las diapositivas plantilla (las que tienen el marcador)
    template_indices = get_template_slide_indices(prs, MARKER)

    if not template_indices:
        print(f"⚠️ No se encontró el marcador {MARKER} en ninguna diapositiva")
        return

    print(f"📋 Encontradas {len(template_indices)} diapositiva(s) con el marcador")
    print(f"📝 Generando para {len(names)} nombre(s)...")

    # Guardar las slides plantilla originales
    template_slides = [prs.slides[i] for i in template_indices]

    # Generar copias para cada nombre
    for name in names:
        for template_slide in template_slides:
            new_slide = duplicate_slide_with_images(prs, template_slide)
            replaced = replace_marker_in_slide(new_slide, name)
            if not replaced:
                print(f"⚠️ No se pudo reemplazar {MARKER} para: {name}")

    # Eliminar las diapositivas plantilla originales (de mayor a menor índice)
    for i in sorted(template_indices, reverse=True):
        delete_slide(prs, i)

    prs.save(OUTPUT_PPTX)
    print(f"✅ Generado: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
